#!/usr/bin/env python3
"""Build the CRPS deck as an editable PowerPoint (.pptx) with fade transitions.

Mirrors the PDF deck (build_deck.py): 25 slides, 16:9, formal medical-keynote
look. Photos are center-cropped to their frames; every slide gets a fade
transition so the deck animates between slides.
"""
import io
import os
import re
import zipfile

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

UPLOADS = "/root/.claude/uploads/d95bbaf8-f767-53d5-9b8b-0f6c3551b17b"
SCRATCH = "/tmp/claude-0/-home-user-kerem-orders/d95bbaf8-f767-53d5-9b8b-0f6c3551b17b/scratchpad"
PDFIMG = os.path.join(SCRATCH, "pdfimg")
OUT = os.path.join(os.path.dirname(__file__), "CRPS-הילית-פז-לביא.pptx")

NAVY = RGBColor(0x11, 0x1A, 0x2B)
NAVY2 = RGBColor(0x1F, 0x2C, 0x47)
IVORY = RGBColor(0xFA, 0xF7, 0xF1)
INK = RGBColor(0x1C, 0x24, 0x34)
MUTED = RGBColor(0x5D, 0x66, 0x75)
DMUTED = RGBColor(0xA9, 0xB2, 0xC4)
COPPER = RGBColor(0xB4, 0x64, 0x2A)
COPPER2 = RGBColor(0x8A, 0x4A, 0x1D)
GOLD = RGBColor(0xDF, 0xA3, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_LINE = RGBColor(0xD8, 0xD2, 0xC6)

FONT = "Arial"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


def P(idx):
    return os.path.join(PDFIMG, idx + ".jpeg")


def U(name):
    return os.path.join(UPLOADS, name)


def slide(bg=IVORY):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def set_rtl(p):
    p._p.get_or_add_pPr().set("rtl", "1")


def txt(s, x, y, w, h, runs, align=PP_ALIGN.RIGHT, size=16, color=INK,
        bold=False, spacing=1.0, anchor=MSO_ANCHOR.TOP, space_after=6):
    """runs: str, or list of paragraphs; each paragraph str or list of (text, kw) runs."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        set_rtl(p)
        if isinstance(para, str):
            para = [(para, {})]
        for t, kw in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = kw.get("font", FONT)
            f.size = Pt(kw.get("size", size))
            f.bold = kw.get("bold", bold)
            f.color.rgb = kw.get("color", color)
    return box


def crop_to(path, w_emu, h_emu, crop_bottom=0.0, top_bias=0.32, max_px=900, q=70):
    img = Image.open(path)
    img = ImageOps.exif_transpose(img).convert("RGB")
    if crop_bottom:
        W, H = img.size
        img = img.crop((0, 0, W, int(H * (1 - crop_bottom))))
    W, H = img.size
    target = w_emu / h_emu
    cur = W / H
    if cur > target:                       # too wide -> crop sides
        nw = int(H * target)
        x0 = (W - nw) // 2
        img = img.crop((x0, 0, x0 + nw, H))
    else:                                  # too tall -> crop, biased to top
        nh = int(W / target)
        y0 = int((H - nh) * top_bias)
        img = img.crop((0, y0, W, y0 + nh))
    img.thumbnail((max_px, max_px), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=q, optimize=True)
    buf.seek(0)
    return buf


def pic(s, path, x, y, w, h, crop_bottom=0.0, top_bias=0.32, max_px=900):
    return s.shapes.add_picture(crop_to(path, int(w), int(h), crop_bottom, top_bias, max_px),
                                x, y, w, h)


def grid(s, paths, x, y, w, h, cols, gap=Inches(0.12), max_px=640):
    rows = (len(paths) + cols - 1) // cols
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rows - 1)) / rows
    for i, p_ in enumerate(paths):
        r, c = divmod(i, cols)
        # RTL reading order: fill right-to-left
        cx = x + (cols - 1 - c) * (cw + gap)
        cy = y + r * (ch + gap)
        pic(s, p_, int(cx), int(cy), int(cw), int(ch), max_px=max_px)


def make3d(shape, depth):
    """Give a shape a real Office 3D extrusion (oblique top-left camera)."""
    from lxml import etree
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    spPr = shape._element.spPr
    sc = etree.SubElement(spPr, A + "scene3d")
    cam = etree.SubElement(sc, A + "camera"); cam.set("prst", "obliqueTopLeft")
    lr = etree.SubElement(sc, A + "lightRig"); lr.set("rig", "threePt"); lr.set("dir", "t")
    sp = etree.SubElement(spPr, A + "sp3d")
    sp.set("extrusionH", str(depth)); sp.set("prstMaterial", "matte")
    bv = etree.SubElement(sp, A + "bevelT"); bv.set("w", "38100"); bv.set("h", "19050")


def card(s, x, y, w, h, title, body, dark=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.07
    sh.fill.solid()
    sh.fill.fore_color.rgb = NAVY2 if dark else WHITE
    sh.line.color.rgb = GOLD if dark else CARD_LINE
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    make3d(sh, 90000)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_rtl(p)
    r = p.add_run(); r.text = title
    r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = GOLD if dark else COPPER2
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.line_spacing = 1.15
    set_rtl(p2)
    r2 = p2.add_run(); r2.text = body
    r2.font.name = FONT; r2.font.size = Pt(11.5)
    r2.font.color.rgb = DMUTED if dark else MUTED


def chapter(s, label, dark=False):
    txt(s, Inches(0.5), Inches(0.12), Inches(4), Inches(0.35), label,
        align=PP_ALIGN.LEFT, size=9, color=DMUTED if dark else MUTED)


def head(s, title, sub=None, dark=False, top=Inches(0.55)):
    txt(s, Inches(1.2), top, Inches(10.93), Inches(0.9), title,
        align=PP_ALIGN.CENTER, size=30, bold=True, color=IVORY if dark else INK)
    if sub:
        txt(s, Inches(1.7), top + Inches(0.85), Inches(9.93), Inches(0.85), sub,
            align=PP_ALIGN.CENTER, size=13, color=DMUTED if dark else MUTED, spacing=1.2)


# ============================================================== photo pools
PAIN_LEAD = [P(i) for i in ["A23", "A12", "A18", "A35", "A43", "A21"]]
PAIN_WALL = [P(i) for i in ["A01", "A02", "A06", "A08", "A19", "A27"]] + \
    [U(n) for n in ["187dc8c6-IMG20250722WA0098.jpg", "e4fe0174-IMG20250722WA0096.jpg",
                    "e4c31d2e-IMG20250722WA0115.jpg", "9a38cdbf-IMG20250722WA0117.jpg",
                    "b429ab7b-IMG20250722WA0097.jpg", "6ee703d5-IMG20250722WA0120.jpg"]]
HOSP = [P(i) for i in ["A04", "A14", "A22", "A25", "A31", "A40"]]
HOSP_WALL = [P(i) for i in ["A03", "A05", "A13", "A15", "A24", "A41"]] + \
    [U(n) for n in ["0d84f904-IMG20250722WA0067.jpg", "fbf59d46-IMG20250722WA0091.jpg",
                    "582b1db3-IMG20250722WA0113.jpg", "a2cf3253-IMG20250722WA0114.jpg"]]
TREAT = [P(i) for i in ["A37", "A32", "A39", "A44", "B02", "B03"]]
TREAT_WALL = [P(i) for i in ["A07", "A10", "A16", "A17", "A20", "A26", "A34", "A36", "A38", "B06"]] + \
    [U(n) for n in ["90cfb7c7-IMG20250722WA0119.jpg", "e7b993c6-IMG20250722WA0122.jpg"]]
REST = [P(i) for i in ["A09", "A11", "A28", "A30", "A42", "B07"]] + \
    [U(n) for n in ["5b89500d-IMG20250722WA0112.jpg", "eb63f106-IMG20250722WA0116.jpg",
                    "880af9ad-IMG20250722WA0107.jpg"]]
FAMILY = [U("a73f490c-IMG20250722WA0118.jpg"), U("203d0b5b-IMG20250722WA0079.jpg")] + \
    [P(i) for i in ["B09", "B15", "B05", "B10"]]
LIFE = [P(i) for i in ["B01", "B04", "B11", "B12", "B13", "B16"]]
GIVE1 = [P(i) for i in ["B14", "B17", "B18", "B21", "B22", "B24", "B25", "B26", "B27", "B28", "B29", "B30"]]
GIVE2 = [P(i) for i in ["B31", "B32", "B33", "B34", "B35", "B37", "B38", "B39", "B40", "B41", "B42", "B43"]]
GIVE3 = [P(i) for i in ["B44", "B45", "B46", "B47", "B48", "B49", "B50", "B36"]]
CROWN = U("f5dc4285-IMG20250722WA0121.jpg")

# ================================================================ 1 · COVER
s = slide(NAVY)
pic(s, CROWN, Inches(7.4), 0, Inches(5.93), SH, crop_bottom=0.14, top_bias=0.15, max_px=1200)
txt(s, Inches(0.7), Inches(1.0), Inches(6.3), Inches(0.4),
    "מצגת אישית בפני צוות מרפאת הכאב", align=PP_ALIGN.CENTER, size=12, color=GOLD)
txt(s, Inches(0.7), Inches(1.5), Inches(6.3), Inches(1.9),
    [[("CRPS", {"size": 96, "bold": True, "color": GOLD})]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(3.45), Inches(6.3), Inches(0.75),
    "תסמונת כאב אזורי מורכבת", align=PP_ALIGN.CENTER, size=28, bold=True, color=IVORY)
txt(s, Inches(0.7), Inches(4.15), Inches(6.3), Inches(0.4),
    "Complex Regional Pain Syndrome", align=PP_ALIGN.CENTER, size=13, color=DMUTED)
txt(s, Inches(0.7), Inches(4.9), Inches(6.3), Inches(0.5),
    [[("הסיפור של ", {"size": 18, "color": IVORY}),
      ("הילית פז לביא", {"size": 18, "bold": True, "color": GOLD})]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(5.5), Inches(6.3), Inches(0.5),
    "״איך לנצל כאב לעוצמה, כוח ונתינה״", align=PP_ALIGN.CENTER, size=15, color=DMUTED)

# ================================================================ 2 · STORY
s = slide()
chapter(s, "פתיחה")
pic(s, P("B19"), Inches(0.6), Inches(0.75), Inches(4.6), Inches(6.0), top_bias=0.2)
txt(s, Inches(5.6), Inches(1.0), Inches(7.1), Inches(0.4), "פתיחה", size=12, color=COPPER, bold=True)
txt(s, Inches(5.6), Inches(1.4), Inches(7.1), Inches(0.9),
    "״אפצ׳י אחד שינה את חיי״", size=30, bold=True)
txt(s, Inches(5.6), Inches(2.4), Inches(7.1), Inches(3.8), [
    "כך נקראת ההרצאה שאני מעבירה היום ברחבי הארץ — כי לפעמים מספיק רגע אחד קטן, תנועה אחת לא נכונה, כדי שהחיים ישתנו מהקצה אל הקצה.",
    "אני אמא, אשת משפחה, אישה פעילה ומלאת חיים. ויום אחד — בלי אזהרה — הגוף שלי הפך לזירת קרב. כאב צורב, בלתי־פוסק, שאף בדיקה לא הצליחה להסביר ואף משכך לא הצליח להשתיק.",
    [("זו לא הרצאה על מחלה. זה סיפור על החיים שמסביב לה — ועל מה שאתם, הרופאים, יכולים לעשות כדי לשנות אותם.",
      {"bold": True, "color": COPPER2})],
], size=14.5, spacing=1.3, space_after=12)

# ============================================================ 3 · WHAT IS
s = slide()
chapter(s, "רקע קליני")
head(s, "מה זה CRPS?",
     "תסמונת כאב כרונית־נוירופתית, המתפתחת לרוב לאחר חבלה או ניתוח — והכאב בה אינו פרופורציונלי לאירוע המחולל")
cards = [
    ("כאב שורף ומתמשך", "כאב עז, צורב ועמוק בגפה — נמשך הרבה אחרי שהפציעה המקורית החלימה."),
    ("אלודיניה", "מגע קל — סדין, בגד, רוח — נחווה ככאב חד. הגוף מתרגם כל גירוי לאיום."),
    ("שינויים נראים", "נפיחות, שינויי צבע וטמפרטורה בעור, הזעה מוגברת, שינויים בציפורניים ובשיער."),
    ("פגיעה תפקודית", "ירידה בטווח תנועה, חולשה ורעד — עד אובדן שימוש בגפה ופגיעה קשה בתפקוד."),
]
for i, (t, b) in enumerate(cards):
    card(s, Inches(0.6 + (3 - i) * 3.09), Inches(2.75), Inches(2.85), Inches(2.9), t, b)
txt(s, Inches(1.2), Inches(6.15), Inches(10.93), Inches(0.6),
    "ההנחה המקובלת: תגובה מוגזמת של מערכת העצבים — ״אזעקת כאב״ שנתקעה במצב פתוח.",
    align=PP_ALIGN.CENTER, size=12, color=MUTED)

# =========================================================== 4 · BUDAPEST
s = slide()
chapter(s, "רקע קליני")
head(s, "קריטריוני בודפשט",
     "האבחנה קלינית בלבד — אין בדיקת מעבדה או הדמיה שמאשרת CRPS. נדרשים תסמינים בשלוש מתוך ארבע קטגוריות, וסימנים בשתיים לפחות")
cards = [
    ("תחושתי", "היפראלגזיה · אלודיניה — כאב לא פרופורציונלי לגירוי."),
    ("וזומוטורי", "אסימטריה בטמפרטורת העור · שינויי צבע — אדום, סגול, חיוור."),
    ("סודומוטורי / בצקת", "נפיחות · שינויים בהזעה ואסימטריה בין הגפיים."),
    ("מוטורי / טרופי", "ירידה בטווח תנועה · חולשה, רעד · שינויים בעור, שיער וציפורניים."),
]
for i, (t, b) in enumerate(cards):
    card(s, Inches(0.6 + (3 - i) * 3.09), Inches(2.75), Inches(2.85), Inches(2.9), t, b)
txt(s, Inches(1.2), Inches(6.15), Inches(10.93), Inches(0.6),
    "ולכן — החשד הקליני שלכם הוא כלי האבחון החשוב ביותר. כאב ״מוגזם״ ביחס לממצאים הוא לעיתים בדיוק הסימן.",
    align=PP_ALIGN.CENTER, size=12, color=MUTED, bold=True)

# ============================================================ 5 · NUMBERS
s = slide(NAVY)
chapter(s, "בשפה של מספרים", dark=True)
head(s, "כמה חזק הכאב הזה?",
     "מדד הכאב של מקגיל (McGill) מדרג כאבים על סולם של 0–50, לפי דיווחי מטופלים — ככה נראה CRPS לעומת כאבים שכולנו מכירים",
     dark=True)
LADDER = [
    ("שבר בעצם", 17, False),
    ("לידה טבעית (ללא אלחוש)", 35, False),
    ("קטיעת אצבע ללא הרדמה", 40, False),
    ("CRPS — הכאב שאני חיה איתו", 42, True),
]
bar_right = Inches(9.3)          # bars grow right-to-left from here
bar_full = Inches(7.4)
for i, (lab, val, is_crps) in enumerate(LADDER):
    y = Inches(2.35 + i * 0.95)
    h = Inches(0.62 if not is_crps else 0.78)
    # label (right column)
    txt(s, Inches(9.5), y + (Inches(0.05) if not is_crps else Inches(0.1)), Inches(3.35), Inches(0.7),
        lab, size=13 if not is_crps else 14.5, bold=is_crps,
        color=GOLD if is_crps else DMUTED, spacing=1.0)
    # track
    tr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(bar_right - bar_full)), y, bar_full, h)
    tr.adjustments[0] = 0.5
    tr.fill.solid(); tr.fill.fore_color.rgb = NAVY2
    tr.line.color.rgb = RGBColor(0x3A, 0x47, 0x63); tr.line.width = Pt(0.5); tr.shadow.inherit = False
    # fill, anchored to the right edge of the track
    fw = Emu(int(int(bar_full) * val / 50))
    f = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(bar_right - fw)), y, fw, h)
    f.adjustments[0] = 0.5
    f.fill.solid(); f.fill.fore_color.rgb = GOLD if is_crps else RGBColor(0x5D, 0x6A, 0x85)
    f.line.fill.background(); f.shadow.inherit = False
    make3d(f, 340000 if is_crps else 200000)
    # value at the left end
    txt(s, Inches(0.55), y - (Inches(0.02) if not is_crps else Inches(0.06)), Inches(1.0), Inches(0.8),
        [[(str(val), {"size": 22 if not is_crps else 30, "bold": True,
                      "color": GOLD if is_crps else DMUTED})]], align=PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(6.35), Inches(10.93), Inches(0.9), [
    [("והנקודה שחשוב לזכור: לידה נגמרת אחרי שעות, שבר מחלים אחרי שבועות — ", {"color": DMUTED, "size": 13}),
     ("הכאב של CRPS נמשך יום אחרי יום, שנה אחרי שנה.", {"color": GOLD, "size": 13, "bold": True})],
], align=PP_ALIGN.CENTER, spacing=1.2)

# ======================================================= 6 · PAIN INTRO
s = slide()
chapter(s, "המסע שלי")
grid(s, PAIN_LEAD, Inches(0.6), Inches(0.75), Inches(5.6), Inches(6.0), cols=2)
txt(s, Inches(6.7), Inches(1.2), Inches(6.0), Inches(0.4), "המסע שלי · פרק א׳", size=12, color=COPPER, bold=True)
txt(s, Inches(6.7), Inches(1.6), Inches(6.0), Inches(0.9), "הכאב שלא רואים", size=30, bold=True)
txt(s, Inches(6.7), Inches(2.6), Inches(6.0), Inches(3.5), [
    "אלה לא תמונות קלות — ובכוונה. ככה נראות שעות הלילה, כשהכאב מגיע לשיא ואין תנוחה אחת שמרגיעה אותו.",
    [("בבדיקות הכל ״תקין״. בעיניים — הכל שבור. זה הפער שכל מטופלת CRPS חיה בתוכו יום־יום.",
      {"bold": True, "color": COPPER2})],
], size=14.5, spacing=1.3, space_after=12)

# ========================================================= 7 · PAIN WALL
s = slide()
chapter(s, "המסע שלי")
head(s, "רגעים שאין להם פילטר",
     "תיעוד אמיתי, לאורך חודשים — כי חשוב לי שתראו את מה שלא נכנס לתיק הרפואי")
grid(s, PAIN_WALL, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=6)

# ========================================================== 8 · HOSPITAL
s = slide()
chapter(s, "המסע שלי")
head(s, "ימים במחלקה",
     "אשפוזים חוזרים, לילות של ניטור, אינספור בדיקות — הקירות הלבנים נהיו בית שני")
grid(s, HOSP, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=3)

# ===================================================== 9 · HOSPITAL WALL
s = slide()
chapter(s, "המסע שלי")
head(s, "בין עירוי לעירוי",
     "גם במחלקה מוצאים רגעים של חיוך — בעיקר כשהילדים מגיעים לבקר")
grid(s, HOSP_WALL, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=5)

# ======================================================== 10 · TREATMENTS
s = slide()
chapter(s, "המסע שלי")
head(s, "הטיפולים",
     "עירויים חוזרים, משאבות, אלקטרודות, חסימות עצביות — מפסיקים לנסות רק מי שמוותרים, ואני לא")
grid(s, TREAT, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=3)

# =================================================== 11 · TREATMENT WALL
s = slide()
chapter(s, "המסע שלי")
head(s, "הגוף כזירת טיפול",
     "צנתרים, מוניטורים ומשאבות — המחיר היומיומי של המרדף אחרי הקלה")
grid(s, TREAT_WALL, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=6)

# ============================================================= 12 · QUOTE
s = slide(NAVY)
pic(s, P("A28"), 0, 0, SW, SH, max_px=1300)
veil = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
veil.fill.solid(); veil.fill.fore_color.rgb = NAVY
veil.line.fill.background(); veil.shadow.inherit = False
# 72% opaque veil so the photo dims behind the quote
sp = veil.fill.fore_color._xFill.find(
    '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
if sp is not None:
    from lxml import etree
    a = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
    a.set('val', '72000')
txt(s, Inches(1.7), Inches(2.35), Inches(9.93), Inches(2.0), [
    "הכאב לא תמיד נראה מבחוץ.",
    "אבל הוא אמיתי — בכל שנייה, בכל נשימה.",
], align=PP_ALIGN.CENTER, size=32, bold=True, color=WHITE, spacing=1.25)
txt(s, Inches(1.7), Inches(4.6), Inches(9.93), Inches(0.5),
    "— מה שכל מטופלת CRPS הייתה רוצה שתדעו", align=PP_ALIGN.CENTER, size=14, color=GOLD)

# ============================================================ 13 · IMPACT
s = slide()
chapter(s, "ההשפעה")
head(s, "מה הכאב לוקח מהחיים", "CRPS הוא לא ״כאב בגפה״ — הוא מערכת חיים שלמה שמתערערת")
cards = [
    ("שינה", "הכאב מתעצם בלילה. שנה שלמה בלי לילה רצוף אחד — וכל יום מתחיל בגירעון."),
    ("תפקוד", "להתלבש, לבשל, לנהוג — פעולות שהיו אוטומטיות הופכות למשימות מתוכננות."),
    ("משפחה", "הילדים לומדים לזהות את ״הימים הקשים״. התפקידים בבית מתהפכים."),
    ("זהות", "המאבק הכי גדול: להישאר אני — אישה, אמא, יוצרת — ולא ״החולה מחדר 4״."),
]
for i, (t, b) in enumerate(cards):
    card(s, Inches(0.6 + (3 - i) * 3.09), Inches(2.75), Inches(2.85), Inches(2.9), t, b)
txt(s, Inches(1.2), Inches(6.15), Inches(10.93), Inches(0.6),
    "ולמרות כל זה — השקפים הבאים הם הסיבה האמיתית שאני כאן.",
    align=PP_ALIGN.CENTER, size=12, color=MUTED, bold=True)

# =============================================================== 14 · REST
s = slide()
chapter(s, "ההשפעה")
head(s, "גם מנוחה היא מאבק",
     "אחרי לילה של כאב, הגוף גובה את שלו — ברכב, על הספה, בכל רגע שמתאפשר")
grid(s, REST, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=5)

# ======================================================== 15 · OUT ANYWAY
s = slide()
chapter(s, "נקודת מפנה")
pic(s, P("A29"), Inches(0.6), Inches(0.75), Inches(3.1), Inches(6.0), top_bias=0.25)
pic(s, P("A33"), Inches(3.85), Inches(0.75), Inches(3.1), Inches(6.0), top_bias=0.25)
txt(s, Inches(7.4), Inches(1.4), Inches(5.3), Inches(0.4), "נקודת מפנה", size=12, color=COPPER, bold=True)
txt(s, Inches(7.4), Inches(1.8), Inches(5.3), Inches(0.9), "ובכל זאת — יוצאים", size=30, bold=True)
txt(s, Inches(7.4), Inches(2.8), Inches(5.3), Inches(3.0), [
    "היד חבושה, הכאב נוכח — ובכל זאת: קניון, הליכה, החיים. הבן שלי לצידי, תומך בכל צעד.",
    [("למדתי שהכאב אולי קובע את הקצב — אבל אני קובעת את הכיוון.", {"bold": True, "color": COPPER2})],
], size=14.5, spacing=1.3, space_after=12)

# ============================================================ 16 · FAMILY
s = slide()
chapter(s, "הכוח")
head(s, "המשפחה — התרופה שאין לה מרשם",
     "הם היו שם בכל אשפוז, בכל לילה לבן. החיבוק שלהם שווה יותר מכל משכך")
grid(s, FAMILY, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=3)

# ============================================================== 17 · LIFE
s = slide()
chapter(s, "הכוח")
head(s, "ממשיכים לחיות — בקול רם",
     "מרוצים, אירועים, חברים — כל תמונה כאן היא החלטה מודעת לא לתת לכאב את המילה האחרונה")
grid(s, LIFE, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=3)

# ====================================================== 18 · WONDER WOMAN
s = slide()
chapter(s, "הכוח")
pic(s, P("B08"), Inches(0.6), 0, Inches(5.2), SH, top_bias=0.12, max_px=1100)
txt(s, Inches(6.3), Inches(1.5), Inches(6.4), Inches(0.4), "הכוח", size=12, color=COPPER, bold=True)
txt(s, Inches(6.3), Inches(1.9), Inches(6.4), Inches(1.6),
    ["לפעמים צריך תחפושת.", "בדרך כלל מספיק אומץ."], size=28, bold=True, spacing=1.15)
txt(s, Inches(6.3), Inches(3.7), Inches(6.4), Inches(2.6), [
    "יום אחד הפכתי את עצמי לוונדר וומן — והבנתי שזו בכלל לא תחפושת. כל מי שקמה בבוקר עם CRPS ובוחרת לחיות את היום הזה בכל זאת — היא גיבורת־על אמיתית.",
    [("וכשצריך כתר — שמים כתר. 👑", {"bold": True, "color": COPPER2})],
], size=14.5, spacing=1.3, space_after=12)

# ======================================================= 19 · GIVE INTRO
s = slide()
chapter(s, "הכוח · נתינה")
grid(s, [P("B20"), P("B23"), P("B19")], Inches(0.6), Inches(1.4), Inches(5.3), Inches(4.7), cols=3)
txt(s, Inches(0.6), Inches(6.2), Inches(5.3), Inches(0.7),
    "ערכה לכבוד אביתר דוד ששב מהשבי · נעלי ״עם ישראל חי״ בעבודת יד · בדרך למסירה",
    align=PP_ALIGN.CENTER, size=10.5, color=MUTED)
txt(s, Inches(6.4), Inches(0.9), Inches(6.3), Inches(0.4), "הכוח · נתינה", size=12, color=COPPER, bold=True)
txt(s, Inches(6.4), Inches(1.3), Inches(6.3), Inches(0.9), "הפכתי את הכאב לנתינה", size=28, bold=True)
txt(s, Inches(6.4), Inches(2.25), Inches(6.3), Inches(4.4), [
    "במקום לשאול ״למה זה קרה לי״, התחלתי לשאול ״מה אני עושה עם זה״. כך נולד מיזם הנתינה שלי: נעליים בעבודת יד עם הכיתוב ״עם ישראל חי״, שאני מעניקה באופן אישי לחטופים ששבו מהשבי ולשורדי הטבח במסיבת הנובה.",
    "והדלת שלי לא נסגרת: אני מארחת אצלי בבית, סביב שולחן וארוחה חמה, שבים מהשבי, שורדי המסיבה והורים שכולים של נרצחי השבעה באוקטובר — ביד פתוחה ובלב שלם.",
    [("כאב מזהה כאב — והנתינה הזאת מרפאת גם אותי.", {"bold": True, "color": COPPER2})],
], size=13, spacing=1.25, space_after=10)

# ======================================================= 20-22 · GIVE WALLS
s = slide()
chapter(s, "הכוח")
head(s, "״עם ישראל חי״ — לשבים מהשבי ולשורדי הנובה",
     "בתמונות: מי שנחטפו לעזה וחזרו הביתה, ומי ששרדו את הטבח במסיבת הנובה בשבעה באוקטובר — כל אחד מקבל זוג מצויר ביד")
grid(s, GIVE1, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=6)

s = slide()
chapter(s, "הכוח")
head(s, "We Will Dance Again — עוד נרקוד",
     "המוטו של שורדי הנובה הפך אצלי לשליחות: להגיע לכל שורד ושורדת, ולהזכיר שהריקוד שלהם עוד לפניהם")
grid(s, GIVE2, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=6)

s = slide()
chapter(s, "הכוח")
head(s, "ובית שהדלת שלו תמיד פתוחה",
     "את חלקם אירחתי אצלי בבית — שבים מהשבי, שורדי המסיבה והורים שכולים של נרצחי השבעה באוקטובר — סביב שולחן, ארוחה חמה ולב פתוח")
grid(s, GIVE3, Inches(0.7), Inches(2.15), Inches(11.93), Inches(4.9), cols=4)

# ========================================================== 23 · LECTURE
s = slide()
chapter(s, "השמעת קול")
pic(s, P("B51"), Inches(1.0), Inches(0.6), Inches(3.6), Inches(6.3), top_bias=0.5, max_px=1000)
txt(s, Inches(5.3), Inches(1.2), Inches(7.4), Inches(0.4), "הכוח · השמעת קול", size=12, color=COPPER, bold=True)
txt(s, Inches(5.3), Inches(1.6), Inches(7.4), Inches(0.9), "היום אני עומדת על במות", size=28, bold=True)
txt(s, Inches(5.3), Inches(2.6), Inches(7.4), Inches(3.6), [
    "אני מרצה לקהילות ברחבי הארץ על החיים עם כאב כרוני: ״אפצ׳י אחד שינה את חיי — איך לנצל כאב לעוצמה, כוח ונתינה״.",
    "מי שפעם לא הצליחה לקום מהמיטה — מחזיקה היום מיקרופון מול קהל. זה לא קרה למרות הכאב. זה קרה בגללו — כי בחרתי לתת לו משמעות.",
    [("והמצגת הזאת, מולכם — היא הבמה הכי חשובה שלי עד היום.", {"bold": True, "color": COPPER2})],
], size=14.5, spacing=1.3, space_after=12)

# ============================================================== 24 · ASKS
s = slide(NAVY)
chapter(s, "המסר", dark=True)
head(s, "מה מטופלת CRPS מבקשת מהרופאים שלה", dark=True)
asks = [
    ("תאמינו לנו.", "הכאב אמיתי גם כשההדמיה תקינה. ״אני מאמין לך״ שווה יותר מעשר בדיקות."),
    ("חשדו מוקדם.", "ב‑CRPS הזמן הוא פרוגנוזה: אבחון והתערבות מוקדמים משנים את מהלך המחלה."),
    ("טפלו רב־תחומית.", "כאב, פיזיותרפיה, תמיכה נפשית ומשפחה — אף טיפול בודד לא מספיק לבדו."),
    ("שמרו על רצף.", "אל תתנו לנו להתחיל כל פעם מאפס מול רופא חדש. ליווי קבוע הוא עוגן."),
    ("תראו את האדם.", "לא רק את הגפה. תמכו בזהות, בעשייה ובתקווה — הן חלק מהטיפול."),
]
for i, (t, b) in enumerate(asks):
    y = Inches(1.75 + i * 1.02)
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10.33), Inches(0.9))
    r.adjustments[0] = 0.16
    r.fill.solid(); r.fill.fore_color.rgb = NAVY2
    r.line.color.rgb = RGBColor(0x3A, 0x47, 0x63); r.line.width = Pt(0.75); r.shadow.inherit = False
    make3d(r, 90000)
    txt(s, Inches(10.6), y + Inches(0.12), Inches(1.0), Inches(0.7),
        [[(str(i + 1), {"size": 26, "bold": True, "color": GOLD})]], align=PP_ALIGN.CENTER)
    txt(s, Inches(1.8), y + Inches(0.17), Inches(8.7), Inches(0.6),
        [[(t + " ", {"size": 14, "bold": True, "color": GOLD}),
          (b, {"size": 13, "color": DMUTED})]], spacing=1.1)

# =========================================================== 25 · CLOSING
s = slide(NAVY)
pic(s, CROWN, Inches(7.4), 0, Inches(5.93), SH, crop_bottom=0.14, top_bias=0.15, max_px=1200)
txt(s, Inches(0.7), Inches(1.7), Inches(6.3), Inches(0.4),
    "תודה שהקשבתם", align=PP_ALIGN.CENTER, size=13, color=GOLD)
txt(s, Inches(0.7), Inches(2.2), Inches(6.3), Inches(1.5),
    "כל יום הוא ניצחון", align=PP_ALIGN.CENTER, size=44, bold=True, color=IVORY)
txt(s, Inches(0.7), Inches(3.8), Inches(6.3), Inches(0.5),
    "הילית פז לביא", align=PP_ALIGN.CENTER, size=20, bold=True, color=GOLD)
txt(s, Inches(0.7), Inches(4.4), Inches(6.3), Inches(0.8),
    "״הכאב הוא חלק מהסיפור שלי — אבל הוא לא כל הסיפור.״",
    align=PP_ALIGN.CENTER, size=15, color=DMUTED)
txt(s, Inches(0.7), Inches(5.3), Inches(6.3), Inches(0.5),
    "#CRPS_Awareness", align=PP_ALIGN.CENTER, size=14, bold=True, color=GOLD)

prs.save(OUT)

# ------------------------------------------------ add fade transitions
TRANS = ('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
         'spd="med"><p:fade/></p:transition>')
tmp = OUT + ".tmp"
with zipfile.ZipFile(OUT) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.namelist():
        data = zin.read(item)
        if re.match(r"ppt/slides/slide\d+\.xml$", item):
            x = data.decode("utf-8")
            if "<p:transition" not in x:
                x = x.replace("</p:sld>", TRANS + "</p:sld>")
            data = x.encode("utf-8")
        zout.writestr(item, data)
os.replace(tmp, OUT)
print(f"Wrote {OUT} ({os.path.getsize(OUT)/1024/1024:.2f} MB), {len(prs.slides.__iter__.__self__._sldIdLst)} slides, fade transitions added")
